"""
Generador de presentacions PowerPoint.
Crea el fitxer .pptx amb l'estil MENAG.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import Optional
import re


# Funció helper per crear colors RGB
def RgbColor(r, g, b):
    """Crea un color RGB compatible amb python-pptx."""
    return RGBColor(r, g, b)

from processors.content_processor import PresentationPlan, SlideContent
from templates.style_config import StyleConfig


def create_presentation(plan: PresentationPlan, output_path: str | Path) -> Path:
    """
    Crea una presentació PowerPoint a partir del pla.

    Args:
        plan: Pla de presentació amb totes les diapositives.
        output_path: Ruta on guardar el fitxer .pptx.

    Returns:
        Path al fitxer creat.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Crear presentació
    prs = Presentation()

    # Configurar mida 16:9
    prs.slide_width = Emu(12192000)  # 13.333 inches
    prs.slide_height = Emu(6858000)  # 7.5 inches

    # Generar cada diapositiva
    for slide_content in plan.slides:
        if slide_content.slide_type == "title":
            _create_title_slide(prs, slide_content, plan)
        elif slide_content.slide_type == "index":
            _create_index_slide(prs, slide_content)
        else:
            _create_content_slide(prs, slide_content)

    # Guardar
    prs.save(str(output_path))
    print(f"Presentació guardada: {output_path}")
    return output_path


def _add_bottom_stripe(slide):
    """Afegeix la franja inferior taronja/marró."""
    # Franja principal
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.1),
        Inches(13.333),
        Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(180, 120, 60)  # Marró/taronja
    shape.line.fill.background()

    # Accent superior (línia taronja més clara)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.05),
        Inches(13.333),
        Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = StyleConfig.ORANGE_PRIMARY
    line.line.fill.background()


def _add_separator_line(slide, top: float = 1.3):
    """Afegeix línia separadora sota el títol."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(top),
        Inches(12.333),
        Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(139, 90, 43)  # Marró
    line.line.fill.background()


def _create_title_slide(prs: Presentation, content: SlideContent, plan: PresentationPlan):
    """Crea la diapositiva de títol/portada."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Títol petit "Capítol X:"
    chapter_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(6), Inches(0.5)
    )
    tf = chapter_label.text_frame
    p = tf.paragraphs[0]
    p.text = f"Capítol {plan.chapter_name.replace('KWC', '')}:"
    p.font.size = Pt(20)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.font.name = StyleConfig.FONT_BODY

    # Títol principal
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0), Inches(7), Inches(2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.title.replace(f"Capítol {plan.chapter_name.replace('KWC', '')}:", "").strip()
    if not p.text:
        p.text = plan.chapter_title
    p.font.size = Pt(40)
    p.font.color.rgb = StyleConfig.ORANGE_PRIMARY
    p.font.name = StyleConfig.FONT_TITLE
    p.font.bold = False

    # Línia separadora
    _add_separator_line(slide, top=4.2)

    # Nom del grup amb membres
    group_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(7), Inches(1.5)
    )
    tf = group_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Processar el nom del grup (pot incloure membres)
    group_text = plan.group_name
    if ":" in group_text:
        # Format "GRUP G: NOM1, NOM2, NOM3"
        p.text = group_text.upper()
    else:
        # Format simple "GRUPG" o "G"
        clean_group = group_text.replace("GRUP", "").strip()
        p.text = f"GRUP {clean_group}"

    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(80, 80, 80)
    p.font.name = StyleConfig.FONT_BODY
    p.font.bold = False

    # Imatge del llibre (si existeix)
    if content.image and content.image.path:
        try:
            slide.shapes.add_picture(
                content.image.path,
                Inches(8), Inches(1.5),
                width=Inches(4)
            )
        except Exception as e:
            print(f"  No s'ha pogut afegir imatge a portada: {e}")

    _add_bottom_stripe(slide)


def _create_index_slide(prs: Presentation, content: SlideContent):
    """Crea la diapositiva d'índex."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Títol
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = content.title
    p.font.size = Pt(40)
    p.font.color.rgb = RgbColor(80, 80, 80)
    p.font.name = StyleConfig.FONT_TITLE

    _add_separator_line(slide)

    # Calcular mida de font segons nombre d'elements
    num_items = len(content.content)
    if num_items <= 6:
        font_size = 18
        space_after = 8
    elif num_items <= 10:
        font_size = 16
        space_after = 6
    else:
        font_size = 14
        space_after = 4

    # Contingut (llista numerada)
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6), Inches(7), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content.content):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        # Netejar número si ja ve amb número
        item_text = item
        if item and item[0].isdigit() and '.' in item[:3]:
            parts = item.split('.', 1)
            if len(parts) > 1:
                item_text = parts[1].strip()

        p.text = f"{i + 1}.{item_text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = StyleConfig.GRAY_DARK
        p.font.name = StyleConfig.FONT_BODY
        p.space_after = Pt(space_after)

    # Imatge decorativa (si existeix)
    if content.image and content.image.path:
        try:
            slide.shapes.add_picture(
                content.image.path,
                Inches(8), Inches(1.8),
                width=Inches(4.5)
            )
        except Exception as e:
            print(f"  No s'ha pogut afegir imatge a índex: {e}")

    _add_bottom_stripe(slide)


def _create_content_slide(prs: Presentation, content: SlideContent):
    """Crea una diapositiva de contingut amb mida de font FIXA i imatge gran."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # MIDA DE FONT FIXA - MAI CANVIA
    FONT_SIZE_NORMAL = 16
    FONT_SIZE_HEADER = 18
    FONT_SIZE_SUBPOINT = 15

    # Títol
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = content.title
    p.font.size = Pt(32)
    p.font.color.rgb = StyleConfig.ORANGE_PRIMARY
    p.font.name = StyleConfig.FONT_TITLE

    _add_separator_line(slide)

    # Determinar layout segons si hi ha imatge (IMATGE GRAN!)
    has_image = content.image and content.image.path
    content_width = Inches(6) if has_image else Inches(12)

    # Contingut amb format FIX
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6), content_width, Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content.content):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        # Analitzar el format de l'ítem
        item_text, _, is_indented, is_header, color, has_bold = _parse_content_item(item)

        # Determinar mida de font FIXA segons tipus
        if is_header:
            font_size = FONT_SIZE_HEADER
        elif is_indented:
            font_size = FONT_SIZE_SUBPOINT
        else:
            font_size = FONT_SIZE_NORMAL

        # Afegir text (amb suport per negreta i subratllat)
        if has_bold:
            _add_text_with_bold(p, item_text, font_size, color, StyleConfig.FONT_BODY)
        else:
            p.text = item_text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = StyleConfig.FONT_BODY

        # Espai FIX entre elements
        if item_text == "" or item_text.strip() == "":
            p.space_after = Pt(2)
        elif is_header:
            p.space_after = Pt(4)
            p.font.bold = True
        else:
            p.space_after = Pt(5)

        # Indentar si és subpunt
        if is_indented:
            p.level = 1

    # Imatge GRAN (si existeix)
    if has_image:
        try:
            img_left = Inches(6.8)
            img_top = Inches(1.5)
            max_width = Inches(6)  # Imatge més gran!

            slide.shapes.add_picture(
                content.image.path,
                img_left, img_top,
                width=max_width
            )
        except Exception as e:
            print(f"  No s'ha pogut afegir imatge a slide {content.number}: {e}")

    _add_bottom_stripe(slide)

    # Afegir notes de l'orador
    if content.speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content.speaker_notes


def _parse_content_item(item: str) -> tuple:
    """
    Analitza un ítem de contingut per determinar el seu format.

    Returns:
        tuple: (text, font_size, is_indented, is_header, color, has_formatting)
    """
    # Colors per defecte
    color_dark = StyleConfig.GRAY_DARK
    color_orange = StyleConfig.ORANGE_PRIMARY

    # Valors per defecte - MIDA FIXA
    font_size = 16  # Mida base fixa
    is_indented = False
    is_header = False
    color = color_dark
    has_formatting = ("**" in item or "__" in item) if item else False

    # Detectar línia buida
    if not item or item.strip() == "":
        return "", 16, False, False, color_dark, False

    stripped = item.strip()

    # Detectar si és indentat (comença amb espais)
    if item.startswith("  ") or item.startswith("\t"):
        is_indented = True

    # Detectar headers/subtítols en **MAJUSCULES** o MAJUSCULES sol
    clean_stripped = stripped.replace("**", "").replace("__", "")
    if clean_stripped.isupper() and len(clean_stripped) < 50 and clean_stripped and not clean_stripped[0].isdigit():
        is_header = True
        color = color_orange

    # Detectar bullet points amb guió o punt
    if stripped.startswith("-") or stripped.startswith("•"):
        is_indented = True

    # Detectar subpunts (comencen amb espai + guió)
    if item.startswith("  -") or item.startswith("    -"):
        is_indented = True

    return item, font_size, is_indented, is_header, color, has_formatting


def _add_text_with_bold(paragraph, text: str, base_font_size: int, base_color, font_name: str):
    """
    Afegeix text a un paràgraf processant **negreta** i __subratllat__.

    Args:
        paragraph: Paràgraf de python-pptx
        text: Text amb possible sintaxi **negreta** o __subratllat__
        base_font_size: Mida de font base
        base_color: Color base
        font_name: Nom de la font
    """
    has_bold = "**" in text
    has_underline = "__" in text

    # Si no hi ha formatació especial, afegir text normal
    if not has_bold and not has_underline:
        paragraph.text = text
        paragraph.font.size = Pt(base_font_size)
        paragraph.font.color.rgb = base_color
        paragraph.font.name = font_name
        return

    # Netejar el text inicial del paràgraf
    paragraph.text = ""

    # Processar text amb negreta i subratllat
    # Primer substituir __text__ per marcadors temporals
    import uuid
    underline_marker = f"__UL_{uuid.uuid4().hex[:4]}__"
    underline_parts = {}

    if has_underline:
        underline_pattern = r'__([^_]+)__'
        underline_matches = re.findall(underline_pattern, text)
        for i, match in enumerate(underline_matches):
            marker = f"{underline_marker}{i}"
            text = text.replace(f"__{match}__", marker, 1)
            underline_parts[marker] = match

    # Ara processar negreta
    pattern = r'\*\*([^*]+)\*\*'
    parts = re.split(pattern, text)

    for i, part in enumerate(parts):
        if not part:
            continue

        # Comprovar si conté marcadors de subratllat
        for marker, underline_text in underline_parts.items():
            if marker in part:
                # Dividir per el marcador
                before, after = part.split(marker, 1)

                if before:
                    run = paragraph.add_run()
                    run.text = before
                    run.font.size = Pt(base_font_size)
                    run.font.color.rgb = base_color
                    run.font.name = font_name
                    if i % 2 == 1:
                        run.font.bold = True

                # Text subratllat
                run = paragraph.add_run()
                run.text = underline_text
                run.font.size = Pt(base_font_size)
                run.font.color.rgb = base_color
                run.font.name = font_name
                run.font.underline = True
                if i % 2 == 1:
                    run.font.bold = True

                if after:
                    run = paragraph.add_run()
                    run.text = after
                    run.font.size = Pt(base_font_size)
                    run.font.color.rgb = base_color
                    run.font.name = font_name
                    if i % 2 == 1:
                        run.font.bold = True

                part = ""  # Part ja processada
                break

        if part:  # Si encara queda part sense processar
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(base_font_size)
            run.font.color.rgb = base_color
            run.font.name = font_name

            # Els índexs senars són el contingut dins **...**
            if i % 2 == 1:
                run.font.bold = True
